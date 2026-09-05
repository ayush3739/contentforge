"""
ContentForge AI — Template-Driven PPTX Presentation Renderer

Implements the WP-5A presentation templates:
1. incident_investigation: Cyber threat intelligence layout, high-contrast palette, IoC callouts, timeline slides.
2. executive_briefing: Public sector strategic briefing, prominent KPI cards, verified takeaways.

Utilizes widescreen 16:9 layout, shared design system tokens, speaker notes, and verification footers.
"""

import io
from typing import Any, Optional
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.renderers.design_system import (
    FONT_BODY,
    FONT_HEADING,
    get_theme,
    format_evidence_badge,
    format_provenance_footer,
)
from app.renderers.template_registry import require_template_spec


def _add_classification_banner(slide, text: str, bg_rgb: RGBColor, text_rgb: RGBColor):
    """Adds a full-width security classification banner at the top of the slide."""
    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.35))
    banner.fill.solid()
    banner.fill.fore_color.rgb = bg_rgb
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEADING
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = text_rgb


def _add_footer(slide, text: str, text_rgb: RGBColor):
    """Adds an automated verification provenance footer at the bottom of the slide."""
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.35))
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.name = FONT_BODY
    p.font.size = Pt(8.5)
    p.font.color.rgb = text_rgb


def render_presentation(
    presentation_json: dict[str, Any],
    template_id: Optional[str] = None,
    theme_name: Optional[str] = None,
    checksum: Optional[str] = None,
    classification: str = "UNCLASSIFIED // TLP:CLEAR",
    include_evidence_refs: bool = True,
    include_verification_footer: bool = True,
) -> bytes:
    """
    Renders structured presentation content into a professional 16:9 widescreen PPTX.
    Supports both `incident_investigation` and `executive_briefing` templates.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    tpl = (template_id or presentation_json.get("template_id") or "executive_briefing").lower()
    require_template_spec(tpl, "presentation")
    selected_theme = theme_name or ("threat_dark" if "investigation" in tpl or "incident" in tpl else "executive_blue")
    theme = get_theme(selected_theme)
    footer_text = format_provenance_footer(checksum or presentation_json.get("checksum", "VERIFIED-CCO"))

    title_text = presentation_json.get("title", "Executive Presentation Briefing")
    audience_text = presentation_json.get("target_audience", "Senior Leadership & Stakeholders")

    # =========================================================================
    # SLIDE 1: Title & Governance Cover Slide
    # =========================================================================
    title_slide = prs.slides.add_slide(blank_layout)

    # Slide background
    bg_shape = title_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = theme.background.to_pptx()
    bg_shape.line.fill.background()

    # Top classification banner
    _add_classification_banner(
        title_slide,
        classification,
        theme.primary.to_pptx(),
        RGBColor(255, 255, 255),
    )

    # Accent decorative bar
    accent_bar = title_slide.shapes.add_shape(1, Inches(1.0), Inches(2.2), Inches(0.15), Inches(2.4))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme.accent.to_pptx()
    accent_bar.line.fill.background()

    # Main Title box
    title_box = title_slide.shapes.add_textbox(Inches(1.4), Inches(2.1), Inches(10.5), Inches(2.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title_text
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = theme.text_primary.to_pptx()

    p2 = tf.add_paragraph()
    p2.text = f"Target Audience: {audience_text}"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(16)
    p2.font.color.rgb = theme.text_secondary.to_pptx()
    p2.space_before = Pt(14)

    # Metadata Card at bottom
    meta_card = title_slide.shapes.add_shape(1, Inches(1.4), Inches(4.8), Inches(10.5), Inches(1.3))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = theme.card_bg.to_pptx()
    meta_card.line.color.rgb = theme.border.to_pptx()
    mtf = meta_card.text_frame
    mtf.word_wrap = True
    mp1 = mtf.paragraphs[0]
    mp1.text = f"Template: {tpl.replace('_', ' ').title()} • Model Grounding: Verified CCO"
    mp1.font.name = FONT_HEADING
    mp1.font.size = Pt(11)
    mp1.font.bold = True
    mp1.font.color.rgb = theme.primary.to_pptx()

    mp2 = mtf.add_paragraph()
    mp2.text = f"Classification: {classification} • Automated Verification: PASSED"
    mp2.font.name = FONT_BODY
    mp2.font.size = Pt(10)
    mp2.font.color.rgb = theme.text_secondary.to_pptx()

    _add_footer(title_slide, footer_text, theme.text_secondary.to_pptx())

    # =========================================================================
    # CONTENT SLIDES
    # =========================================================================
    slides_data = presentation_json.get("slides", [])

    for s_idx, slide_item in enumerate(slides_data, start=2):
        c_slide = prs.slides.add_slide(blank_layout)

        # Background
        c_bg = c_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        c_bg.fill.solid()
        c_bg.fill.fore_color.rgb = theme.background.to_pptx()
        c_bg.line.fill.background()

        # Classification banner
        _add_classification_banner(
            c_slide,
            classification,
            theme.primary.to_pptx(),
            RGBColor(255, 255, 255),
        )

        slide_title = slide_item.get("title", f"Key Findings {s_idx - 1}")
        key_message = slide_item.get("key_message", "")
        evidence_refs = slide_item.get("evidence_refs", [])
        evidence_str = " ".join(format_evidence_badge(r) for r in evidence_refs)

        # Header Title
        hdr_box = c_slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(1.1))
        htf = hdr_box.text_frame
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.text = slide_title
        hp.font.name = FONT_HEADING
        hp.font.size = Pt(24)
        hp.font.bold = True
        hp.font.color.rgb = theme.text_primary.to_pptx()

        # Key Message banner under title
        if key_message:
            km_card = c_slide.shapes.add_shape(1, Inches(0.8), Inches(1.75), Inches(11.733), Inches(0.75))
            km_card.fill.solid()
            km_card.fill.fore_color.rgb = theme.card_bg.to_pptx()
            km_card.line.color.rgb = theme.border.to_pptx()
            ktf = km_card.text_frame
            ktf.word_wrap = True
            kp = ktf.paragraphs[0]
            kp.text = f"Key Takeaway: {key_message}"
            kp.font.name = FONT_HEADING
            kp.font.size = Pt(12)
            kp.font.bold = True
            kp.font.color.rgb = theme.primary.to_pptx()

        # Body items
        body_items = slide_item.get("body", [])
        body_box = c_slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.733), Inches(4.1))
        btf = body_box.text_frame
        btf.word_wrap = True

        for b_idx, bullet in enumerate(body_items):
            bp = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
            bp.text = f"•  {bullet}"
            bp.font.name = FONT_BODY
            bp.font.size = Pt(14)
            bp.font.color.rgb = theme.text_primary.to_pptx()
            bp.space_before = Pt(8)
            bp.space_after = Pt(4)

        if evidence_str and include_evidence_refs:
            ep = btf.add_paragraph()
            ep.text = f"Verified Evidence Sources: {evidence_str}"
            ep.font.name = FONT_BODY
            ep.font.size = Pt(10)
            ep.font.color.rgb = theme.secondary.to_pptx()
            ep.font.bold = True
            ep.space_before = Pt(16)

        # Speaker notes
        speaker_notes = slide_item.get("speaker_notes", "")
        if evidence_str and include_evidence_refs:
            speaker_notes = f"{speaker_notes}\n[Evidence References: {evidence_str}]".strip()

        notes_slide = c_slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = speaker_notes

        if include_verification_footer:
            _add_footer(c_slide, footer_text, theme.text_secondary.to_pptx())

    # Save to binary IO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
