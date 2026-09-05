"""
ContentForge AI — Wave 3 Automated Tests
Tests the 6 controlled templates across PPTX, DOCX, and Infographic SVG,
the visual design tokens, structured verification findings, and claim lineage.
"""

import io
import json
import pytest
from pptx import Presentation
from docx import Document as DocxDocument

from app.renderers.design_system import (
    get_theme,
    format_evidence_badge,
    format_provenance_footer,
    CLASSIFICATION_BANNERS,
    SEVERITY_COLORS,
)
from app.renderers.template_registry import (
    TEMPLATE_SPECS,
    ArtifactTemplateConfig,
    get_template_spec,
    get_default_template_id,
    list_templates_for_type,
)
from app.renderers.pptx_renderer import render_presentation
from app.renderers.docx_renderer import render_document
from app.renderers.infographic_renderer import render_infographic_svg
from app.ai.verification.verifier import verify_artifact, VerificationReport, VerificationIssue


def test_design_system_tokens():
    """Validates visual tokens, palettes, and provenance formatters."""
    exec_theme = get_theme("executive_blue")
    assert exec_theme.name == "executive_blue"
    assert exec_theme.primary.hex.startswith("#")
    assert exec_theme.background.to_pptx() is not None
    assert exec_theme.text_primary.to_docx() is not None

    threat_theme = get_theme("threat_dark")
    assert threat_theme.name == "threat_dark"
    assert threat_theme.is_dark is True

    minimal_theme = get_theme("modern_minimal")
    assert minimal_theme.name == "modern_minimal"

    # Evidence badge formatter
    assert format_evidence_badge("E-01") == "[E-01]"
    assert format_evidence_badge("[E-02]") == "[E-02]"

    # Provenance footer formatter
    footer = format_provenance_footer("a1b2c3d4e5f6")
    assert "ContentForge AI" in footer
    assert "a1b2c3d4" in footer

    # Classification banners
    assert "UNCLASSIFIED // TLP:CLEAR" in CLASSIFICATION_BANNERS.values()
    assert "CONFIDENTIAL // INTERNAL USE ONLY" in CLASSIFICATION_BANNERS.values()


def test_template_registry_specs():
    """Validates the 6 demo template specifications and configuration contracts."""
    expected_ids = [
        "incident_investigation",
        "executive_briefing",
        "executive_summary",
        "security_advisory",
        "incident_brief",
        "executive_snapshot",
    ]
    for tid in expected_ids:
        spec = get_template_spec(tid)
        assert spec is not None, f"Missing template spec for {tid}"
        assert spec.name
        assert len(spec.features) > 0

    # Test filtering by artifact type
    pptx_specs = list_templates_for_type("presentation")
    assert len(pptx_specs) == 2
    assert {s.id for s in pptx_specs} == {"incident_investigation", "executive_briefing"}

    docx_specs = list_templates_for_type("executive_summary") + list_templates_for_type("advisory")
    assert len(docx_specs) == 2

    svg_specs = list_templates_for_type("infographic")
    assert len(svg_specs) == 2
    assert {s.id for s in svg_specs} == {"incident_brief", "executive_snapshot"}

    # Default template mapping
    assert get_default_template_id("presentation") == "executive_briefing"
    assert get_default_template_id("executive_summary") == "executive_summary"
    assert get_default_template_id("advisory") == "security_advisory"
    assert get_default_template_id("infographic") == "executive_snapshot"

    # Contract validation
    cfg = ArtifactTemplateConfig(
        artifact_type="presentation",
        template_id="incident_investigation",
        brand_theme="threat_dark",
    )
    assert cfg.template_id == "incident_investigation"
    assert cfg.include_evidence_refs is True


def test_pptx_renderer_templates():
    """Tests rendering both PPTX templates to genuine 16:9 widescreen presentations."""
    sample_presentation = {
        "title": "Operation Nightshade Incident Investigation",
        "target_audience": "Executive Leadership & SecOps",
        "checksum": "e3b0c44298fc1c149afbf4c8996fb92427ae41e4649b934ca495991b7852b855",
        "slides": [
            {
                "slide_number": 1,
                "title": "Incident Timeline & Kill Chain",
                "key_message": "SIEM anomaly detected at T0 across 14 payment gateway clusters.",
                "body": [
                    "00:00 (T0): SIEM outbound beaconing detected from node cluster B-12.",
                    "04:30 (T+4h): Network quarantine enforced on ingress/egress firewalls.",
                    "12:00 (T+12h): Security patch KB-9912 deployed across all nodes.",
                ],
                "speaker_notes": "Emphasize that customer PII was not accessed or exfiltrated [E-01].",
                "evidence_refs": ["E-01", "E-02"],
            },
            {
                "slide_number": 2,
                "title": "Threat Vector & Indicators of Compromise",
                "key_message": "Vulnerability CVE-2024-3094 mitigated with 100% patch coverage.",
                "body": [
                    "Threat actor attempted privilege escalation via rogue PAM library.",
                    "All 14 node hashes validated against official vendor checksums.",
                    "External forensic firm engaged for independent ledger verification.",
                ],
                "speaker_notes": "Highlight zero unauthorized data movement verified by database integrity hashes [E-03].",
                "evidence_refs": ["E-03", "E-04"],
            },
        ],
        "key_takeaways": [
            "14 core database nodes successfully quarantined within response SLA.",
            "Zero unauthorized customer PII exfiltration verified across object storage.",
            "Total incident remediation expenditure capped at $2.5M under corporate insurance.",
        ],
    }

    # 1. Test executive_briefing template
    pptx_bytes_exec = render_presentation(
        sample_presentation,
        template_id="executive_briefing",
        theme_name="executive_blue",
        classification="UNCLASSIFIED // TLP:CLEAR",
    )
    assert pptx_bytes_exec.startswith(b"PK\x03\x04")  # Valid ZIP/PPTX container
    prs_exec = Presentation(io.BytesIO(pptx_bytes_exec))
    assert prs_exec.slide_width.inches == pytest.approx(13.333, 0.01)  # 16:9 widescreen
    assert prs_exec.slide_height.inches == pytest.approx(7.5, 0.01)
    assert len(prs_exec.slides) >= 3  # Cover + content slides + takeaways

    # 2. Test incident_investigation template
    pptx_bytes_inv = render_presentation(
        sample_presentation,
        template_id="incident_investigation",
        theme_name="threat_dark",
        classification="CONFIDENTIAL // INTERNAL USE ONLY",
    )
    assert pptx_bytes_inv.startswith(b"PK\x03\x04")
    prs_inv = Presentation(io.BytesIO(pptx_bytes_inv))
    assert len(prs_inv.slides) >= 3
    # Verify speaker notes are retained
    has_notes = any(s.has_notes_slide and s.notes_slide.notes_text_frame.text for s in prs_inv.slides)
    assert has_notes is True


def test_docx_renderer_templates():
    """Tests rendering both DOCX templates into structured Word documents."""
    sample_doc = {
        "title": "Enterprise Cybersecurity Advisory & Strategy Brief",
        "checksum": "9f86d081884c7d659a2feaa0c55ad015a3bf4f1b2b0b822cd15d6c15b0f00a08",
        "severity": "HIGH",
        "executive_overview": "A coordinated response mitigated remote execution risk across enterprise nodes.",
        "key_findings": [
            "Zero unauthorized customer PII exfiltration detected [E-01].",
            "14 core server instances isolated and patched with KB-9912 [E-02].",
        ],
        "impact": [
            "Financial: Remediation capped at $2.5M under corporate indemnity insurance [E-03].",
            "Operational: Downtime restricted to 42 minutes during initial routing.",
        ],
        "recommended_actions": [
            "Rotate all service credentials and TLS certificates within 24 hours.",
            "Deploy secondary honeypot sensors in edge VPC subnets.",
        ],
        "indicators_of_compromise": [
            {"type": "SHA-256", "value": "e3b0c44298fc1c149afbf4c8996fb92427ae41e4649b934ca495991b7852b855", "context": "Payload"},
            {"type": "IPv4", "value": "198.51.100.42", "context": "C2 Server"},
        ],
        "affected_entities": ["Node Cluster B-12", "Payment Gateway Microservice"],
    }

    # 1. Test executive_summary template
    docx_bytes_exec = render_document(
        sample_doc,
        template_id="executive_summary",
        theme_name="executive_blue",
    )
    assert docx_bytes_exec.startswith(b"PK\x03\x04")
    doc_exec = DocxDocument(io.BytesIO(docx_bytes_exec))
    assert len(doc_exec.tables) >= 1  # Document Control block table
    doc_text_exec = "\n".join(p.text for p in doc_exec.paragraphs)
    assert "Enterprise Cybersecurity Advisory" in doc_text_exec

    # 2. Test security_advisory template
    docx_bytes_adv = render_document(
        sample_doc,
        template_id="security_advisory",
        theme_name="threat_dark",
        classification="RESTRICTED // LAW ENFORCEMENT SENSITIVE",
    )
    assert docx_bytes_adv.startswith(b"PK\x03\x04")
    doc_adv = DocxDocument(io.BytesIO(docx_bytes_adv))
    assert len(doc_adv.tables) >= 2  # Threat banner / Affected systems / IoC table
    doc_text_adv = "\n".join(p.text for p in doc_adv.paragraphs)
    assert "Enterprise Cybersecurity Advisory" in doc_text_adv
    assert "Required Mitigation & Containment Checklist" in doc_text_adv


def test_infographic_svg_templates():
    """Tests server-side rendering of both Infographic templates into vector SVGs."""
    sample_infographic = {
        "title": "Incident Intelligence & Operational Impact",
        "checksum": "5e884898da28047151d0e56f8dc6292773603d0d6aabbdd62a11ef721d1542d8",
        "metrics": [
            {"label": "Systems Quarantined", "value": "14", "evidence_ref": "E-01"},
            {"label": "Remediation Cap", "value": "$2.5M", "evidence_ref": "E-02"},
            {"label": "Data Leakage", "value": "0 KB", "evidence_ref": "E-03"},
            {"label": "Patch SLA", "value": "24h", "evidence_ref": "E-04"},
        ],
        "timeline": [
            {"time": "00:00 (T0)", "event": "SIEM Anomaly Detected", "detail": "Beaconing from 14 nodes.", "status": "critical"},
            {"time": "04:30 (T+4h)", "event": "Cluster Quarantine", "detail": "Firewall ingress/egress cut.", "status": "warning"},
            {"time": "24:00 (T+24h)", "event": "100% Remediation", "detail": "Patch KB-9912 applied.", "status": "success"},
        ],
        "comparison_bars": [
            {"label": "Patch Verification", "percent": 100, "value": "100%"},
            {"label": "Infrastructure Shielding", "percent": 92, "value": "92%"},
            {"label": "Audit Traceability", "percent": 100, "value": "100%"},
        ],
        "key_takeaways": [
            {"title": "Zero PII Compromised", "text": "Cryptographic proof confirmed zero data leakage across clusters."},
            {"title": "Containment Enforced", "text": "All 14 perimeter nodes quarantined within 4 hours."},
        ],
    }

    # 1. Test executive_snapshot template
    svg_bytes_snapshot = render_infographic_svg(
        sample_infographic,
        template_id="executive_snapshot",
        theme_name="executive_blue",
    )
    svg_str_snapshot = svg_bytes_snapshot.decode("utf-8")
    assert svg_str_snapshot.startswith("<svg")
    assert svg_str_snapshot.strip().endswith("</svg>")
    assert "Incident Intelligence" in svg_str_snapshot
    assert "SYSTEMS QUARANTINED" in svg_str_snapshot
    assert "14" in svg_str_snapshot
    assert "[E-01]" in svg_str_snapshot

    # 2. Test incident_brief template
    svg_bytes_brief = render_infographic_svg(
        sample_infographic,
        template_id="incident_brief",
        theme_name="threat_dark",
        classification="CONFIDENTIAL // INTERNAL USE ONLY",
    )
    svg_str_brief = svg_bytes_brief.decode("utf-8")
    assert svg_str_brief.startswith("<svg")
    assert "Chronology" in svg_str_brief
    assert "SIEM Anomaly Detected" in svg_str_brief
    assert "CONFIDENTIAL // INTERNAL USE ONLY" in svg_str_brief
    assert "ContentForge AI" in svg_str_brief


def test_structured_verification_issues_and_lineage():
    """Validates that verifier.py generates itemized structured issues with suggested fixes."""
    cco = {
        "title": "Incident CCO",
        "claims": [
            {"id": "CLM-001", "text": "14 database server nodes were isolated within 24 hours."},
            {"id": "CLM-002", "text": "Zero customer PII was accessed or exfiltrated."},
        ],
        "numbers": [{"text": "14"}, {"text": "24"}],
    }

    # Artifact containing one cited slide and one slide lacking citations
    artifact = {
        "artifact_type": "presentation",
        "title": "Incident Summary",
        "slides": [
            {
                "title": "Overview",
                "body": ["14 database server nodes isolated within 24 hours."],
                "evidence_refs": ["E-01"],
            },
            {
                "title": "Unverified Financials",
                "body": ["Total unverified loss was 999 million dollars."],
                "evidence_refs": [],  # missing evidence
            },
        ],
    }

    evidence = [
        {"chunk_id": "chunk-001", "text": "14 database server nodes were isolated within 24 hours of alert."},
        {"chunk_id": "chunk-002", "text": "Zero customer PII was accessed or exfiltrated."},
    ]

    report = verify_artifact(
        artifact=artifact,
        cco=cco,
        evidence=evidence,
    )

    assert isinstance(report, VerificationReport)
    assert report.citation_coverage >= 0.0
    assert report.citation_coverage <= 1.0
    assert len(report.issues) > 0

    # Verify structured issues
    for issue in report.issues:
        assert isinstance(issue, VerificationIssue)
        assert issue.id.startswith("ISS-")
        assert issue.severity in ["CRITICAL", "HIGH", "MEDIUM", "LOW"]
        assert issue.suggested_fix
        assert len(issue.suggested_fix) > 0
